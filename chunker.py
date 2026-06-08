from pathlib import Path
import config


def _chunk_text(text: str, source: str) -> list[dict]:
    """将文本按 CHUNK_SIZE/OVERLAP 分块。"""
    chunks = []
    start = 0
    idx = 0
    while start < len(text):
        end = start + config.CHUNK_SIZE
        chunk_text = text[start:end].strip()
        if chunk_text:
            chunks.append({
                "id": f"{source}#{idx}",
                "text": chunk_text,
                "metadata": {"source": source, "page": 0},
            })
            idx += 1
        start += config.CHUNK_SIZE - config.CHUNK_OVERLAP
    return chunks


def _load_txt(file_path: str) -> str:
    with open(file_path, encoding="utf-8") as f:
        return f.read()


def _load_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _load_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _load_md(file_path: str) -> str:
    return _load_txt(file_path)


def _load_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        parts.append(f"表名：{sheet.title}")
        headers = [str(h) if h is not None else "" for h in rows[0]]
        for row in rows[1:]:
            cells = [str(v) if v is not None else "" for v in row]
            if any(cells):
                parts.append(" | ".join(
                    f"{h}: {v}" for h, v in zip(headers, cells) if h or v
                ))
        parts.append("")
    return "\n".join(parts)


def _load_xls(file_path: str) -> str:
    import xlrd
    wb = xlrd.open_workbook(file_path)
    parts = []
    for sheet in wb.sheets():
        if sheet.nrows == 0:
            continue
        parts.append(f"表名：{sheet.name}")
        headers = [str(sheet.cell_value(0, c)) for c in range(sheet.ncols)]
        for r in range(1, sheet.nrows):
            cells = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
            if any(cells):
                parts.append(" | ".join(
                    f"{h}: {v}" for h, v in zip(headers, cells) if h or v
                ))
        parts.append("")
    return "\n".join(parts)


def _load_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            ph = getattr(shape, "placeholder_format", None)
            # placeholder type 15=TITLE, 3=CENTER_TITLE
            if ph is not None and ph.type in (15, 3):
                title_text = text
            else:
                body_texts.append(text)
        header = f"--- 第{i}页：{title_text} ---" if title_text else f"--- 第{i}页 ---"
        parts.append(header)
        parts.extend(body_texts)
        parts.append("")
    return "\n".join(parts)


def _load_csv(file_path: str) -> str:
    import csv
    parts = []
    rows = []
    for encoding in ("utf-8-sig", "gbk", "latin-1"):
        try:
            with open(file_path, encoding=encoding, errors="replace", newline="") as f:
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=",\t")
                except csv.Error:
                    dialect = csv.excel
                rows = list(csv.reader(f, dialect))
            break
        except Exception:
            continue
    if not rows:
        return ""
    headers = [str(h) for h in rows[0]]
    for row in rows[1:]:
        if any(v.strip() for v in row):
            parts.append(" | ".join(
                f"{h}: {v}" for h, v in zip(headers, row) if h or v
            ))
    return "\n".join(parts)


def _load_html(file_path: str) -> str:
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        SKIP_TAGS = {"script", "style", "head"}
        BLOCK_TAGS = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "div", "li", "br", "tr"}

        def __init__(self):
            super().__init__()
            self._parts = []
            self._skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in self.SKIP_TAGS:
                self._skip += 1
            if tag in self.BLOCK_TAGS and self._skip == 0:
                self._parts.append("\n")

        def handle_endtag(self, tag):
            if tag in self.SKIP_TAGS:
                self._skip = max(0, self._skip - 1)

        def handle_data(self, data):
            if self._skip == 0:
                self._parts.append(data)

        def get_text(self):
            return "".join(self._parts)

    html_content = ""
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            with open(file_path, encoding=encoding, errors="replace") as f:
                html_content = f.read()
            break
        except Exception:
            continue

    parser = _TextExtractor()
    parser.feed(html_content)
    return parser.get_text()


def load_and_chunk(file_path: str) -> list[dict]:
    """解析文档并分块。支持 .pdf/.docx/.txt/.md/.xlsx/.xls/.pptx/.csv/.html/.htm。"""
    path = Path(file_path)
    suffix = path.suffix.lower()
    source = path.name

    loaders = {
        ".txt":  _load_txt,
        ".md":   _load_md,
        ".pdf":  _load_pdf,
        ".docx": _load_docx,
        ".xlsx": _load_xlsx,
        ".xls":  _load_xls,
        ".pptx": _load_pptx,
        ".csv":  _load_csv,
        ".html": _load_html,
        ".htm":  _load_html,
    }

    if suffix not in loaders:
        raise ValueError(
            f"不支持的文件格式: {suffix}（支持：{', '.join(sorted(loaders))}）"
        )

    text = loaders[suffix](file_path)
    return _chunk_text(text, source)
