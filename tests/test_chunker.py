import pytest
from pathlib import Path
import chunker
import config


def test_chunk_text_basic():
    """短文本应产生 1 个 chunk。"""
    chunks = chunker._chunk_text("Python 是一种编程语言。", source="test.txt")
    assert len(chunks) == 1
    assert chunks[0]["text"] == "Python 是一种编程语言。"
    assert chunks[0]["metadata"]["source"] == "test.txt"


def test_chunk_text_splits_long_text(monkeypatch):
    """超过 CHUNK_SIZE 的文本应拆分为多个 chunk。"""
    monkeypatch.setattr(config, "CHUNK_SIZE", 10)
    monkeypatch.setattr(config, "CHUNK_OVERLAP", 2)
    text = "A" * 25
    chunks = chunker._chunk_text(text, source="long.txt")
    assert len(chunks) > 1
    assert all(c["metadata"]["source"] == "long.txt" for c in chunks)


def test_chunk_ids_are_unique():
    text = "X" * 600
    chunks = chunker._chunk_text(text, source="dup.txt")
    ids = [c["id"] for c in chunks]
    assert len(ids) == len(set(ids))


def test_load_and_chunk_txt(tmp_path):
    txt_file = tmp_path / "sample.txt"
    txt_file.write_text("这是第一段内容。\n这是第二段内容。", encoding="utf-8")
    chunks = chunker.load_and_chunk(str(txt_file))
    assert len(chunks) >= 1
    assert chunks[0]["metadata"]["source"] == "sample.txt"
    full_text = " ".join(c["text"] for c in chunks)
    assert "第一段" in full_text


def test_load_and_chunk_unsupported_format(tmp_path):
    bad_file = tmp_path / "file.xyz"
    bad_file.write_text("content")
    with pytest.raises(ValueError, match="不支持的文件格式"):
        chunker.load_and_chunk(str(bad_file))


def test_load_md(tmp_path):
    path = tmp_path / "readme.md"
    path.write_text("# 标题\n\n正文内容。", encoding="utf-8")
    text = chunker._load_md(str(path))
    assert "标题" in text
    assert "正文内容" in text


def test_load_xlsx(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "销售"
    ws.append(["产品", "数量"])
    ws.append(["苹果", 10])
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    text = chunker._load_xlsx(str(path))
    assert "表名：销售" in text
    assert "产品: 苹果" in text


def test_load_xls(tmp_path):
    xlwt = pytest.importorskip("xlwt")
    wb = xlwt.Workbook()
    ws = wb.add_sheet("库存")
    ws.write(0, 0, "仓库")
    ws.write(0, 1, "数量")
    ws.write(1, 0, "北京")
    ws.write(1, 1, "500")
    path = tmp_path / "test.xls"
    wb.save(str(path))
    text = chunker._load_xls(str(path))
    assert "表名：库存" in text
    assert "仓库: 北京" in text


def test_load_pptx(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "季度总结"
    slide.placeholders[1].text = "营收增长15%"
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    text = chunker._load_pptx(str(path))
    assert "季度总结" in text
    assert "营收增长" in text


def test_load_csv(tmp_path):
    path = tmp_path / "data.csv"
    path.write_text("姓名,部门,薪资\n张三,研发,20000\n李四,市场,18000", encoding="utf-8")
    text = chunker._load_csv(str(path))
    assert "姓名: 张三" in text
    assert "部门: 研发" in text


def test_load_html(tmp_path):
    path = tmp_path / "page.html"
    path.write_text(
        "<html><head><title>Test</title></head><body>"
        "<h1>标题</h1><p>正文内容</p>"
        "<script>alert('skip')</script>"
        "</body></html>",
        encoding="utf-8",
    )
    text = chunker._load_html(str(path))
    assert "标题" in text
    assert "正文内容" in text
    assert "skip" not in text


def test_load_and_chunk_new_formats(tmp_path):
    """load_and_chunk dispatches correctly for .md and .csv."""
    md = tmp_path / "note.md"
    md.write_text("## 笔记\n内容", encoding="utf-8")
    chunks = chunker.load_and_chunk(str(md))
    assert len(chunks) >= 1
    assert chunks[0]["metadata"]["source"] == "note.md"

    csv_file = tmp_path / "data.csv"
    csv_file.write_text("key,val\nfoo,bar", encoding="utf-8")
    chunks2 = chunker.load_and_chunk(str(csv_file))
    assert len(chunks2) >= 1
